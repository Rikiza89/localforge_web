"""
document_extractor.py — テキスト以外のドキュメント形式からプレーンテキストを抽出する。
PDF / DOCX / XLSX / PPTX に対応。各パッケージは optional — 未インストール時は空文字を返す。
"""
from __future__ import annotations
import logging
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)


def extract_text(path: Path) -> str:
    """
    ファイル形式を判別してプレーンテキストを返す。
    抽出失敗・パッケージ未インストール時は空文字を返す（クラッシュしない）。
    """
    suffix = path.suffix.lower()
    try:
        if suffix == ".pdf":
            return _extract_pdf(path)
        elif suffix in (".docx", ".doc"):
            return _extract_docx(path)
        elif suffix in (".xlsx", ".xls"):
            return _extract_xlsx(path)
        elif suffix in (".pptx", ".ppt"):
            return _extract_pptx(path)
        else:
            return path.read_text(encoding="utf-8", errors="replace")
    except Exception as exc:
        logger.warning("ドキュメント抽出エラー (%s): %s", path.name, exc)
        return ""


def _extract_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        logger.debug("pypdf 未インストール — PDF 抽出スキップ: %s", path.name)
        return ""
    reader = PdfReader(str(path))
    pages = []
    for page in reader.pages:
        text = page.extract_text() or ""
        pages.append(text)
    return "\n".join(pages)


def _extract_docx(path: Path) -> str:
    try:
        import docx
    except ImportError:
        logger.debug("python-docx 未インストール — DOCX 抽出スキップ: %s", path.name)
        return ""
    doc = docx.Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # テーブル内のテキストも取得
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    paragraphs.append(cell.text.strip())
    return "\n".join(paragraphs)


def _extract_xlsx(path: Path) -> str:
    try:
        import openpyxl
    except ImportError:
        logger.debug("openpyxl 未インストール — XLSX 抽出スキップ: %s", path.name)
        return ""
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    lines = []
    for sheet in wb.worksheets:
        lines.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join(str(c) for c in row if c is not None)
            if row_text.strip():
                lines.append(row_text)
    return "\n".join(lines)


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        logger.debug("python-pptx 未インストール — PPTX 抽出スキップ: %s", path.name)
        return ""
    prs = Presentation(str(path))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"[Slide {i}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
    return "\n".join(lines)
